# ============================================================
# 🎯 AUTO-GENERATE POWERPOINT: Median OS & PFS Analysis
# ============================================================

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# -------------------------------
# ⚙️ Paths
# -------------------------------
PFS_DIR = r"C:\LARVOL_WORK\Median_PFS\outputs\scatter_PFS"
OS_DIR = r"C:\LARVOL_WORK\Median_PFS\outputs\scatter_OS"
FEATURE_DIR = r"C:\LARVOL_WORK\Median_PFS\outputs\models"
OUTPUT_PPT = r"C:\LARVOL_WORK\Median_PFS\outputs\MedianPFS_OS_Summary.pptx"

# -------------------------------
# 🪧 Helper functions
# -------------------------------
def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_section_title(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(1)).text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.LEFT

def add_image_slide(prs, title, img_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    left = Inches(0.5)
    top = Inches(0.5)
    slide.shapes.add_picture(img_path, left, top + Inches(1), height=Inches(5))
    tf = slide.shapes.add_textbox(left, top, Inches(9), Inches(0.8)).text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.bold = True
    p.font.size = Pt(20)

def add_text_slide(prs, title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# -------------------------------
# 🖼️ Build the presentation
# -------------------------------
prs = Presentation()

# 1️⃣ Title
add_title_slide(
    prs,
    "Clinical Trial Insights: Median OS & PFS Analysis",
    "Correlation and Feature Analysis of ORR, DoR, and Product Effects"
)

# 2️⃣ Summary Slide
summary_text = (
    "This analysis explores clinical endpoints (Median OS and Median PFS)\n"
    "across multiple oncology trials, focusing on:\n\n"
    "• Objective Response Rate (ORR)\n"
    "• Duration of Response (DoR)\n"
    "• Treatment Product comparison\n\n"
    "We generated scatter plots to examine correlations between response\n"
    "metrics and survival outcomes, along with feature importance insights."
)
add_text_slide(prs, "Project Summary", summary_text)

# 3️⃣ Median PFS Scatter Plots
add_section_title(prs, "Median PFS — Key Relationships")
for name in ["1_ORR_vs_MedianPFS.png", "2_DoR_vs_MedianPFS.png", "3_Product_vs_MedianPFS.png"]:
    path = os.path.join(PFS_DIR, name)
    if os.path.exists(path):
        add_image_slide(prs, f"Plot: {name.replace('_', ' ').replace('.png','')}", path)

# 4️⃣ Median OS Scatter Plots
add_section_title(prs, "Median OS — Key Relationships")
for name in ["1_ORR_vs_MedianOS.png", "2_DoR_vs_MedianOS.png", "3_Product_vs_MedianOS.png"]:
    path = os.path.join(OS_DIR, name)
    if os.path.exists(path):
        add_image_slide(prs, f"Plot: {name.replace('_', ' ').replace('.png','')}", path)

# 5️⃣ Feature Importance (optional)
add_section_title(prs, "Model Feature Importance (LightGBM)")
for name in ["feature_importance_grouped.png", "feature_importance_os.png"]:
    path = os.path.join(FEATURE_DIR, name)
    if os.path.exists(path):
        add_image_slide(prs, f"Feature Importance — {os.path.basename(name)}", path)

# 6️⃣ Summary of Insights
insights_text = (
    "📊 Key Observations:\n"
    "• Higher ORR and longer DoR generally correlate with improved PFS and OS.\n"
    "• Some products demonstrate significantly higher median outcomes.\n"
    "• Feature importance highlights ORR and DoR as top predictive variables.\n\n"
    "🧭 Next Steps:\n"
    "• Explore multivariate modeling to isolate effect of product type.\n"
    "• Validate correlations with updated datasets."
)
add_text_slide(prs, "Insights & Next Steps", insights_text)

# Save PowerPoint
prs.save(OUTPUT_PPT)
print(f"✅ PowerPoint saved successfully → {OUTPUT_PPT}")
