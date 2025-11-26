import os
import numpy as np
import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt
from scipy.stats import pearsonr
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# ===============================
# ⚙️ CONFIG
# ===============================
# INPUT_PFS = r"C:\LARVOL_WORK\Median_PFS\data\MedianPFS_cleaned.xlsx"
INPUT_PFS=r"C:\LARVOL_WORK\Median_PFS\data\MedianPFS_PMOA_with_precise_clusters_v3.xlsx"
INPUT_OS  = r"C:\LARVOL_WORK\Median_PFS\data\MedianOS_cleaned.xlsx"
OUTPUT_DIR = r"C:\LARVOL_WORK\Median_PFS\outputs\corr_full_report"
os.makedirs(OUTPUT_DIR, exist_ok=True)

# ===============================
# 📂 LOAD DATA
# ===============================
pfs = pd.read_excel(INPUT_PFS)
osd = pd.read_excel(INPUT_OS)
print(f"✅ Loaded: PFS={pfs.shape}, OS={osd.shape}")

# Convert to numeric
for df in (pfs, osd):
    for col in [
        "Objective Response Rate Percentage",
        "Duration of Response Median",
        "Median PFS",
        "Median OS"
    ]:
        if col in df.columns:
            df[col] = pd.to_numeric(df[col], errors="coerce")

# ===============================
# 🧮 COMPUTE ORR×DoR + log(ORR×DoR)
# ===============================
for df in (pfs, osd):
    df["ORR×DoR"] = (
        df["Objective Response Rate Percentage"] * df["Duration of Response Median"]
    )
    df["log(ORR×DoR)"] = np.log1p(df["ORR×DoR"])

# ===============================
# 📈 CORRELATION + SCATTER PLOT
# ===============================
def corr_and_plot(df, x, y, color, fname):
    sub = df[[x, y]].dropna()
    if len(sub) < 3:
        return [fname, None, None]

    r, p = pearsonr(sub[x], sub[y])
    plt.figure(figsize=(8, 6))
    sns.regplot(
        data=sub, x=x, y=y,
        scatter_kws={"alpha": 0.6, "s": 45, "color": color},
        line_kws={"color": "red", "lw": 2}
    )
    plt.title(f"{x} vs {y}", fontsize=13)
    plt.xlabel(x)
    plt.ylabel(y)
    plt.text(0.02, 0.95, f"r = {r:.3f}\np = {p:.2e}",
             transform=plt.gca().transAxes, ha="left", va="top", fontsize=10,
             bbox=dict(boxstyle="round", facecolor="white", alpha=0.8))
    plt.tight_layout()
    path = os.path.join(OUTPUT_DIR, fname)
    plt.savefig(path, bbox_inches="tight", dpi=300)
    plt.close()
    return [f"{x} vs {y}", round(r, 3), "{:.2e}".format(p), path]

# ===============================
# 🔹 COMPUTE ALL 8 COMBINATIONS
# ===============================
results = []
results.append(corr_and_plot(pfs, "Objective Response Rate Percentage", "Median PFS", "#1f77b4", "orr_vs_medianPFS.png"))
results.append(corr_and_plot(osd, "Objective Response Rate Percentage", "Median OS", "#ff7f0e", "orr_vs_medianOS.png"))
results.append(corr_and_plot(pfs, "Duration of Response Median", "Median PFS", "#2ca02c", "dor_vs_medianPFS.png"))
results.append(corr_and_plot(osd, "Duration of Response Median", "Median OS", "#9467bd", "dor_vs_medianOS.png"))
results.append(corr_and_plot(pfs, "ORR×DoR", "Median PFS", "#17becf", "orrdor_vs_medianPFS.png"))
results.append(corr_and_plot(osd, "ORR×DoR", "Median OS", "#d62728", "orrdor_vs_medianOS.png"))
results.append(corr_and_plot(pfs, "log(ORR×DoR)", "Median PFS", "#8c564b", "log_orrdor_vs_medianPFS.png"))
results.append(corr_and_plot(osd, "log(ORR×DoR)", "Median OS", "#e377c2", "log_orrdor_vs_medianOS.png"))

corr_df = pd.DataFrame(results, columns=["Feature Comparison", "Correlation (r)", "P-value", "Plot Path"])
corr_df.dropna(subset=["Correlation (r)"], inplace=True)
corr_path = os.path.join(OUTPUT_DIR, "correlation_summary.xlsx")
corr_df.to_excel(corr_path, index=False)

# ===============================
# 🧾 POWERPOINT
# ===============================
prs = Presentation()

# --- Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Correlation Analysis — ORR, DoR, ORR×DoR vs PFS/OS"
slide.placeholders[1].text = "Including log-transformed metrics (log(ORR×DoR))"

# --- Schema of Inclusion Slide
schema_slide = prs.slides.add_slide(prs.slide_layouts[5])
schema_slide.shapes.title.text = "Schema of Inclusion"

schema_text = """
Dataset Features:
• Trial ID, Arm ID — trial identifiers
• Type — arm type (Experimental / Control / Comparator)
• Product, Dosage — treatment identifiers
• Arm N — sample size
• Source Name, Publication Year — metadata
• Objective Response Rate N — number of responders
✓ Objective Response Rate Percentage (ORR%)
✓ Duration of Response Median (DoR)
✓ ORR×DoR (combined response metric)
✓ log(ORR×DoR) (log-normalized response metric)
✓ Median PFS (Progression-Free Survival)
✓ Median OS (Overall Survival)
• Response_Count_Duration, Response_Percentage_Duration — engineered variables (not used here)
"""
shape = schema_slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5))
tf = shape.text_frame
p = tf.add_paragraph()
p.text = schema_text
p.font.size = Pt(16)
p.space_after = Pt(5)
p.alignment = PP_ALIGN.LEFT

# --- Correlation Table Slide
slide_tbl = prs.slides.add_slide(prs.slide_layouts[5])
slide_tbl.shapes.title.text = "Correlation Summary Table"
rows, cols = corr_df.shape
table = slide_tbl.shapes.add_table(rows + 1, cols - 1, Inches(0.5), Inches(1.5), Inches(9), Inches(0.5 + rows * 0.25)).table
for j, col in enumerate(corr_df.columns[:-1]):
    table.cell(0, j).text = col
for i in range(rows):
    for j in range(cols - 1):
        table.cell(i + 1, j).text = str(corr_df.iloc[i, j])

# --- One Slide per Plot
for _, row in corr_df.iterrows():
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = row["Feature Comparison"]
    img = row["Plot Path"]
    if os.path.exists(img):
        slide.shapes.add_picture(img, Inches(1), Inches(1.5), width=Inches(8))

# --- Save
ppt_path = os.path.join(OUTPUT_DIR, "Median_PFS_OS_Report.pptx")
prs.save(ppt_path)

print("\n🎯 REPORT COMPLETE:")
print(f"📘 Excel Summary → {corr_path}")
print(f"📊 PowerPoint → {ppt_path}")
